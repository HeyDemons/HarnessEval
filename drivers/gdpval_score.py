#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import os
import re
from pathlib import Path
from typing import Any


def artifact_text(path: Path) -> str:
    suffix = path.suffix.lower()
    try:
        if suffix == ".docx":
            from docx import Document
            return "\n".join(p.text for p in Document(path).paragraphs)
        if suffix == ".pptx":
            from pptx import Presentation
            return "\n".join(shape.text for slide in Presentation(path).slides for shape in slide.shapes if hasattr(shape, "text"))
        if suffix == ".xlsx":
            from openpyxl import load_workbook
            book = load_workbook(path, read_only=True, data_only=True)
            return "\n".join(" | ".join(str(value) for value in row if value is not None) for sheet in book for row in sheet.iter_rows(values_only=True))
        if suffix == ".pdf":
            from pypdf import PdfReader
            return "\n".join(page.extract_text() or "" for page in PdfReader(path).pages)
        if suffix in {".txt", ".md", ".csv", ".json", ".html"}:
            return path.read_text(encoding="utf-8", errors="replace")
    except Exception as exc:
        return f"[extraction failed: {type(exc).__name__}: {exc}]"
    return ""


def rubric_items(raw: Any) -> list[dict[str, Any]]:
    if raw is None:
        raise ValueError("GDPVal rubric is missing")
    if isinstance(raw, str):
        raw = json.loads(raw)
    if isinstance(raw, list):
        items = [item if isinstance(item, dict) else {"criterion": str(item)} for item in raw]
        if not items:
            raise ValueError("GDPVal rubric is empty")
        return items
    if isinstance(raw, dict):
        for key in ("rubrics", "criteria", "items"):
            if isinstance(raw.get(key), list):
                return rubric_items(raw[key])
        return [{"criterion": key, "detail": value} for key, value in raw.items()]
    return [{"criterion": str(raw)}]


def score(workspace: Path, authority: Path) -> dict[str, Any]:
    gold = json.loads((authority / "gold.json").read_text())
    case_path = authority.parent / "input" / "case.json"
    case = json.loads(case_path.read_text()) if case_path.is_file() else {}
    references = {str(name) for name in case.get("reference_files") or []}
    files = sorted(path for path in workspace.rglob("*") if path.is_file() and path.name not in references)
    inventory = []
    for path in files:
        text = artifact_text(path)
        inventory.append({"path": str(path.relative_to(workspace)), "bytes": path.stat().st_size, "text": text[:12000]})
    items = rubric_items(gold.get("rubric_json"))
    prompt = {
        "rubric_items": items,
        "expected_deliverables": gold.get("deliverable_files"),
        "submitted_artifacts": inventory,
        "instructions": (
            "Grade each rubric item against only the submitted artifact evidence. Return JSON with "
            "items [{index, score, max_score, rationale}] and overall_score from 0 to 1. Missing evidence scores zero."
        ),
    }
    # Imported here, not at module scope, for the same reason docx/pptx/openpyxl/pypdf are
    # above: it is needed only when a rubric is actually graded. At module scope it made
    # drivers.gdpval_score unimportable wherever requests is absent -- the run venv on the
    # 4090 among them -- which took test_benchmark_scorers down with it even though the
    # functions it covers do no I/O. Every other caller in this platform uses urllib.
    import requests

    base = os.environ["HARNESS_API_BASE"].rstrip("/")
    endpoint = base if base.endswith("/chat/completions") else base + "/chat/completions"
    response = requests.post(endpoint, headers={"Authorization": f"Bearer {os.environ['HARNESS_API_KEY']}", "Content-Type": "application/json"}, json={
        "model": os.environ["HARNESS_MODEL"], "temperature": 0, "response_format": {"type": "json_object"},
        "messages": [{"role": "system", "content": "You are an independent GDPval rubric grader."}, {"role": "user", "content": json.dumps(prompt, ensure_ascii=False)}],
    }, timeout=float(os.environ.get("HARNESS_API_TIMEOUT_S", "300")))
    response.raise_for_status()
    content = response.json()["choices"][0]["message"]["content"]
    content = re.sub(r"^```(?:json)?\s*|\s*```$", "", content.strip())
    verdict = json.loads(content)
    value = verdict.get("overall_score")
    if not isinstance(value, (int, float)):
        rows = verdict.get("items") or []
        earned = sum(float(row.get("score") or 0) for row in rows)
        possible = sum(float(row.get("max_score") or 0) for row in rows)
        value = earned / possible if possible else None
    return {"authority": "gdpval_independent_model_rubric_proxy", "comparability": "proxy", "score": float(value) if value is not None else None, "artifacts": [item["path"] for item in inventory], "rubric": verdict}


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--workspace", type=Path, required=True)
    parser.add_argument("--authority", type=Path, required=True)
    args = parser.parse_args()
    print(json.dumps(score(args.workspace, args.authority), ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
